"""
Sean Scully, October 2024

"main.py" - csv2pptx.py
program to generate pptxs from simple csv data files
aligns bullet points according to spaces before each word in the cell
generate generic powerpoint we can style afterward.

Goals:
standardize powerpoints
make new powerpoint generation easier
focus more on content, organization, and accessibility than wrangling slides

Future:
Done-add dosis-semibold for font
-test functionality of image slides addition
    -pull title of image slide from a separate csv
    -expand to pull from a folder of image files and a master slide title csv
-expand to read a folder of data files and output multiple pptx
-utilize filename for title slide "title.text ="
-build prompt for autogenerating the csv, with all the commas and spaces, or auto-format sample text
    -intermediary files?

"""

import csv
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# Function to set the font for a text run
def set_font(para, font_name, font_size):
    for run in para.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = True  # Make font bold if necessary

# Function to add a bullet slide
def add_bullet_slide(prs, title_text, bullet_points):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title_text
    set_font(title_shape.text_frame.paragraphs[0], 'Dosis SemiBold', 44)  # Set title font ##

    # Access the text frame of the body shape
    tf = body_shape.text_frame

    # Clear existing text in the text frame
    tf.clear()

    # Add bullet points
    for bullet in bullet_points:
        level = bullet.count(" ")  # Count leading spaces to determine bullet level
        text = bullet.strip()  # Remove leading spaces for the bullet text

        # Create a new paragraph for each bullet
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.level = max(0, level - 2)  # Use level - 1 to correct the bullet level
        set_font(p, 'Dosis SemiBold', 32)  # Set bullet point font #

"""#perhaps process folder of images with pillow to standardize size (and save modified file with standard filename) before we import?

# Function to add a slide with a title and two side-by-side images
def add_images_slide(prs, slide_title, image_path1, image_path2):
    # Choose a blank layout or any other layout which accommodates images
    slide_layout = prs.slide_layouts[5]  # Using layout index 5 for a blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Set the title
    title_shape = slide.shapes.title
    title_shape.text = slide_title

    # Define dimensions and positions for the images
    left1 = Inches(1)  # Adjust as necessary
    top1 = Inches(1.5) # Adjust as necessary
    left2 = Inches(5)  # Adjust as necessary
    top2 = Inches(1.5) # Align the second image with the first

    # Set the width and height for the images
    width = Inches(3)  # Adjust size of the images
    height = Inches(4) # Adjust size of the images

    # Add images to the slide
    slide.shapes.add_picture(image_path1, left1, top1, width, height)
    slide.shapes.add_picture(image_path2, left2, top2, width, height)

"""

# Create a presentation object
prs = Presentation()

# Add the title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Make With SOLIDWORKS!"
set_font(title.text_frame.paragraphs[0], 'Dosis SemiBold', 48)  # Set title font for subtitle
subtitle.text = "TTU Library Emerging Tech Workshop\nInstructor Sean Scully\n\n2025" # can add new line characters inside the strings 
                # applies to data pulled from the csv?
set_font(subtitle.text_frame.paragraphs[0], 'Dosis SemiBold', 36)  # Set font for the subtitle


# Read data from CSV and create bullet slides
with open('data.csv', newline='', encoding='utf-8') as csvfile:
    csv_reader = csv.reader(csvfile)
    for row in csv_reader:
        title_text = row[0]  # The first column as slide title
        bullet_points = row[1:]  # The rest as bullet points
        add_bullet_slide(prs, title_text, bullet_points)

# Insert Dosis-semibold font? large point? https://python-pptx.readthedocs.io/en/latest/user/text.html

#font = run.font
#font.name = 'Dosis-SemiBold'



"""#simply copy paste for each image slide? or add "images" folder and "imageSlideTitles.csv" and run through a loop?

# Add a slide with title and two images at the end
image_title = "Images Side by Side"
image_path1 = "image1.png"  # Replace with your actual path
image_path2 = "image2.png"  # Replace with your actual path
add_images_slide(prs, image_title, image_path1, image_path2)
"""

# Save the presentation
prs.save('test.pptx')

"""
Instead of directly assigning p.level = level, we adjust it with p.level = max(0, level - 1). This means:

    If the bullet point has no leading spaces (indicating it's a top-level bullet), level would be 0, so after subtracting 1, it remains 0.
    For sub-bullets with one leading space, this results in p.level = 0 for the first level, p.level = 1 for the second level, and so forth, correctly starting bullet levels from 0 instead of 1.
"""


"""import csv
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

# Function to add a bullet slide
def add_bullet_slide(prs, title_text, bullet_points):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title_text

    # Access the text frame of the body shape
    tf = body_shape.text_frame

    # Clear existing text in the text frame
    tf.clear()

    # Add bullet points
    for bullet in bullet_points:
        level = bullet.count(" ")  # Count leading spaces to determine bullet level
        text = bullet.strip()  # Remove leading spaces for the bullet text

        # Create a new paragraph for each bullet
        p = tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.level = level  # Set the level based on leading spaces

# Create a presentation object
prs = Presentation()

# Add the title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

# Read data from CSV and create bullet slides
with open('data.csv', newline='', encoding='utf-8') as csvfile:
    csv_reader = csv.reader(csvfile)
    for row in csv_reader:
        title_text = row[0]  # The first column as slide title
        bullet_points = row[1:]  # The rest as bullet points
        add_bullet_slide(prs, title_text, bullet_points)

# Save the presentation
prs.save('test.pptx')"""
"""


    CSV Format: Sub-bullets are indicated by leading spaces in the CSV. You can adjust how many spaces represent different bullet levels based on your preference.

    Determining Levels: The bullet level is determined by counting leading spaces in bullet.count(" "). The bullet text is extracted by stripping these spaces using bullet.strip().

    Simplified Bullet Handling: This approach reduces the need for complex parsing and more easily aligns with how you might visually represent content in a slide.

Benefits of This Approach:

    Simplicity: It’s easier to modify the CSV by simply adding spaces for sub-bullets, making it more visually understandable.
    Maintainability: There's less code complexity in determining levels, which aids in debugging.
    Flexibility: You can easily adjust the presentation structure just by modifying the CSV without changing code logic.

This should help avoid issues with data parsing and bullet point creation.
"""
