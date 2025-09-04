import os
from pptx import Presentation

current_dir = os.path.dirname(os.path.realpath(__file__))

prs = Presentation(current_dir + '\\template_sample.pptx')
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello World"
subtitle.text = "python-pptx was here!?"

prs.save('test-template.pptx')